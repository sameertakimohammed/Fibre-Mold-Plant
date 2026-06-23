"""Multi-month PowerPoint deck: a custom range produces per-month detail slides
plus cross-month trend slides; a single month omits the trend slides."""
import io

from pptx import Presentation


def _slide_titles(content: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(content))
    titles = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)
    return titles


def _seed_two_months(client, headers):
    shifts = [
        {"work_date": "2024-01-10", "shift": "Day", "qty": 1000, "fuel_use": 200,
         "downtime_min": 60, "sched_hours": 8, "comment": "mold change"},
        {"work_date": "2024-01-20", "shift": "Day", "qty": 1100, "fuel_use": 210,
         "downtime_min": 30, "sched_hours": 8, "comment": "washing molds"},
        {"work_date": "2024-02-12", "shift": "Day", "qty": 1300, "fuel_use": 220,
         "downtime_min": 45, "sched_hours": 8, "comment": "pump repair"},
    ]
    for s in shifts:
        r = client.post("/api/v1/shifts", headers=headers, json=s)
        assert r.status_code in (201, 409), r.text


def test_pptx_multi_month_has_trend_slides(client, admin_headers):
    _seed_two_months(client, admin_headers)
    r = client.get("/api/v1/reports/report.pptx", headers=admin_headers,
                   params={"start": "2024-01-01", "end": "2024-02-29", "period": "Monthly"})
    assert r.status_code == 200, r.text
    assert "presentationml" in r.headers["content-type"]

    titles = _slide_titles(r.content)
    # title + (2 months x 3 detail) + 3 cross-month summary slides = 10
    assert len(titles) == 10, titles
    joined = " | ".join(titles)
    assert "January 2024 — Production Performance" in titles
    assert "February 2024 — Diesel Fuel Consumption" in titles
    assert "Monthly Trend Comparison" in joined
    assert "Downtime Category Comparison" in joined
    assert "Observations & Trends" in joined


def test_pptx_single_month_omits_trend_slides(client, admin_headers):
    _seed_two_months(client, admin_headers)
    r = client.get("/api/v1/reports/report.pptx", headers=admin_headers,
                   params={"start": "2024-01-01", "end": "2024-01-31", "period": "Monthly"})
    assert r.status_code == 200, r.text
    titles = _slide_titles(r.content)
    # title + 3 detail slides for the single month, no trend slides
    assert len(titles) == 4, titles
    assert "Monthly Trend Comparison" not in " | ".join(titles)
