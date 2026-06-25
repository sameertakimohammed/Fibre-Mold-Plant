"""PowerPoint (.pptx) report builder (python-pptx).

A management-facing deck for the monthly report over a CUSTOM DATE RANGE that
can span several months. The deck gives:

  * a title slide for the period,
  * for EACH month in the range, three detail slides — Production Performance
    (daily output chart + KPI tiles), Downtime Analysis (chart + table) and
    Diesel Consumption (chart),
  * when the range covers 2+ months, cross-month SUMMARY slides — a Monthly
    Trend Comparison, a Downtime Category Comparison, and an auto-generated
    Observations & Trends slide so trends over the period are easy to spot.

Figures come from collect_trend_data() — soft-deleted rows excluded, consistent
with the dashboard KPIs and the xlsx/pdf builders.

build_report_pptx(db, start, end[, period_label]) -> (pptx_bytes, filename)
"""
import io
from datetime import date

from sqlalchemy.orm import Session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from .report_data import report_filename
from .report_trends import collect_trend_data, deck_narrative, CAUSE_ORDER, TrendData, MonthData

AMBER = RGBColor(0xF5, 0xA6, 0x23)
INK = RGBColor(0x14, 0x19, 0x20)      # near-black slide background
PANEL = RGBColor(0x1E, 0x26, 0x30)    # tile fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTE = RGBColor(0x9A, 0xA4, 0xB2)
GREEN = RGBColor(0x36, 0xB3, 0x7E)
RED = RGBColor(0xE5, 0x6A, 0x4E)
PURPLE = RGBColor(0x9B, 0x7E, 0xDE)
BLUE = RGBColor(0x4E, 0x9A, 0xE5)

# Chart series palette (applied to the first N series).
SERIES_COLORS = [AMBER, BLUE, GREEN, PURPLE, RED]

SW, SH = Inches(13.333), Inches(7.5)


def _fmt(n) -> str:
    if isinstance(n, float) and not float(n).is_integer():
        return f"{n:,.1f}"
    return f"{int(round(n)):,}"


# --------------------------------------------------------------------------- #
# Primitives
# --------------------------------------------------------------------------- #
def _bg(slide, color=INK):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = color
    return tb


def _accent_bar(slide, left, top, width=Inches(2.2), color=AMBER):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    return bar


def _caption(slide, text, top=Inches(6.95), color=MUTE, size=11, prefix="✦ "):
    """A slim full-width caption band at the bottom of a slide (used for the
    AI-written one-line summaries)."""
    if not text:
        return None
    _text(slide, Inches(0.85), top, Inches(11.6), Inches(0.5),
          [(prefix + text, size, False, color)])


def _slide_header(prs, title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_bar(s, Inches(0.85), Inches(0.45), width=Inches(0.7))
    _text(s, Inches(0.85), Inches(0.55), Inches(11.6), Inches(0.7),
          [(title, 24, True, WHITE)])
    if sub:
        _text(s, Inches(0.87), Inches(1.2), Inches(11.6), Inches(0.4),
              [(sub, 13, False, MUTE)])
    return s


def _tile(slide, left, top, w, h, value, label, accent=AMBER, vsize=26, lsize=10):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = PANEL; box.line.width = Pt(0.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.size = Pt(vsize); r.font.bold = True; r.font.name = "Calibri"; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(lsize); r2.font.name = "Calibri"; r2.font.color.rgb = MUTE
    _accent_bar(slide, left + Inches(0.12), top + Inches(0.1),
                width=w - Inches(0.24), color=accent)
    return box


def _set_cell(cell, text, *, bold=False, size=12, color=WHITE,
              align=PP_ALIGN.LEFT, fill=PANEL):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    cell.margin_left = Pt(6); cell.margin_right = Pt(6)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = "Calibri"; r.font.color.rgb = color


def _table(slide, left, top, width, headers, rows, col_widths=None,
           body_size=12, head_size=12):
    """headers: list[str]; rows: list[list[str]]. First column left-aligned,
    the rest right-aligned. Returns the table."""
    n_cols = len(headers)
    shape = slide.shapes.add_table(len(rows) + 1, n_cols, left, top, width,
                                   Inches(0.4) * (len(rows) + 1))
    table = shape.table
    if col_widths:
        for c, w in enumerate(col_widths):
            table.columns[c].width = w
    for c, h in enumerate(headers):
        _set_cell(table.cell(0, c), h, bold=True, color=INK, fill=AMBER, size=head_size,
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
    for ri, row in enumerate(rows, 1):
        fill = INK if ri % 2 else PANEL
        for c, val in enumerate(row):
            _set_cell(table.cell(ri, c), val, fill=fill, size=body_size,
                      align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT,
                      bold=(c == 0))
    return table


def _chart(slide, chart_type, left, top, width, height, categories, series,
           *, legend=False, number_format=None, cat_font=9):
    """series: list[(name, [values], color|None)]."""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values, _color in series:
        cd.add_series(name, values, number_format=number_format)
    gframe = slide.shapes.add_chart(chart_type, left, top, width, height, cd)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    # Colour the series.
    for i, plot_series in enumerate(chart.series):
        _, _, color = series[i]
        if color is not None:
            fmt = plot_series.format
            fmt.fill.solid(); fmt.fill.fore_color.rgb = color
            if chart_type == XL_CHART_TYPE.LINE:
                plot_series.format.line.color.rgb = color
    # Axis label sizing so dense day-of-month categories stay readable.
    try:
        chart.category_axis.tick_labels.font.size = Pt(cat_font)
        chart.value_axis.tick_labels.font.size = Pt(9)
    except Exception:  # pragma: no cover - some chart types lack an axis
        pass
    return chart


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def _title_slide(prs, trend: TrendData):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _accent_bar(s, Inches(0.9), Inches(2.55), width=Inches(2.4))
    _text(s, Inches(0.85), Inches(2.7), Inches(11.5), Inches(1.2),
          [("Fibre Mold Plant", 44, True, WHITE)])
    _text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.7),
          [("Production Report", 24, False, AMBER)])
    span = trend.span
    if trend.multi:
        span = f"{span}  ·  {len(trend.months)} months"
    _text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.6),
          [(f"Period: {span}", 16, False, MUTE)])
    _text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
          [("Golden Manufacturers · Recycling Department", 12, False, MUTE)])
    return s


def _month_production_slide(prs, month: MonthData, targets: dict, ai_note=None):
    s = _slide_header(prs, f"{month.label} — Production Performance",
                      "Daily forming machine output")
    m = month.metrics
    # Daily output chart (left).
    cats = [d["day"] for d in month.by_day] or ["—"]
    vals = [round(d["qty"]) for d in month.by_day] or [0]
    _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.85), Inches(1.7),
           Inches(7.6), Inches(4.9), cats, [("Trays", vals, AMBER)],
           cat_font=8)

    # KPI tiles (right column).
    tx, tw, th, gap = Inches(8.75), Inches(3.7), Inches(0.74), Inches(0.14)
    top = Inches(1.7)
    # Monthly production targets (supervisor's plan) shown on the product tiles.
    mix30 = month.mix["30's trays"]
    mix12 = month.mix["12's cartons"]
    t30 = targets.get("prod_30")
    t12 = targets.get("prod_12")
    lbl30 = f"30's · {round(mix30 / t30 * 100)}% of {_fmt(t30)} target" if t30 else "30's trays"
    lbl12 = f"12's · {round(mix12 / t12 * 100)}% of {_fmt(t12)} target" if t12 else "12's cartons"
    tiles = [
        (_fmt(m["total_qty"]), "Total trays produced", AMBER),
        (_fmt(mix30), lbl30, GREEN),
        (_fmt(mix12), lbl12, BLUE),
        (_fmt(month.mix["Hot pressed"]), "Hot pressed", PURPLE),
        (_fmt(month.mix["Labelled"]), "Labelled", AMBER),
        (_fmt(m["avg_per_day"]), "Avg trays / day", GREEN),
    ]
    for i, (v, l, acc) in enumerate(tiles):
        _tile(s, tx, top + i * (th + gap), tw, th, v, l, acc, vsize=18, lsize=9)
    _caption(s, ai_note)
    return s


def _month_downtime_slide(prs, month: MonthData, ai_note=None):
    s = _slide_header(prs, f"{month.label} — Downtime Analysis by Category",
                      f"Total downtime {_fmt(month.metrics['total_downtime_hrs'])} hrs "
                      f"· {_fmt(month.metrics['downtime_pct'])}% of scheduled")
    causes = month.causes
    ordered = [(c, causes.get(c, 0)) for c in CAUSE_ORDER if causes.get(c, 0) > 0]
    total = sum(v for _, v in ordered) or 1
    if ordered:
        cats = [c for c, _ in ordered]
        hrs = [round(v / 60, 1) for _, v in ordered]
        _chart(s, XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.85), Inches(1.7),
               Inches(6.7), Inches(4.8), cats, [("Hours lost", hrs, RED)],
               cat_font=10)
        rows = [[c, _fmt(round(v / 60, 1)), f"{round(v / total * 100)}%"]
                for c, v in ordered]
        _table(s, Inches(7.8), Inches(1.9), Inches(4.65),
               ["Category", "Hours", "Share"], rows,
               col_widths=[Inches(2.75), Inches(1.05), Inches(0.85)])
    else:
        _text(s, Inches(0.9), Inches(3.2), Inches(11), Inches(1),
              [("No downtime recorded this month.", 16, False, MUTE)])
    _caption(s, ai_note)
    return s


def _month_fuel_slide(prs, month: MonthData, ai_note=None):
    m = month.metrics
    s = _slide_header(prs, f"{month.label} — Diesel Fuel Consumption",
                      f"{_fmt(m['total_fuel'])} L burned · {_fmt(m['fuel_eff'])} L per 1,000 trays")
    cats = [d["day"] for d in month.by_day] or ["—"]
    vals = [round(d["fuel"]) for d in month.by_day] or [0]
    _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.85), Inches(1.7),
           Inches(8.2), Inches(4.9), cats, [("Diesel (L)", vals, BLUE)],
           cat_font=8)
    tx, tw, th = Inches(9.35), Inches(3.1), Inches(1.2)
    _tile(s, tx, Inches(1.8), tw, th, _fmt(m["total_fuel"]), "Total diesel (L)", BLUE)
    _tile(s, tx, Inches(3.2), tw, th, _fmt(m["fuel_eff"]), "L / 1,000 trays", AMBER)
    _tile(s, tx, Inches(4.6), tw, th, _fmt(m["avg_speed"]), "Avg speed (pcs/hr)", GREEN)
    _caption(s, ai_note)
    return s


def _trend_comparison_slide(prs, trend: TrendData, overall_summary=None):
    s = _slide_header(prs, f"Monthly Trend Comparison",
                      f"{trend.span} · how the plant is tracking month over month")
    labels = [m.label.split(" ")[0][:3] + " " + m.label.split(" ")[1][-2:] for m in trend.months]
    # Output column chart.
    qty = [m.metrics["total_qty"] for m in trend.months]
    _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.85), Inches(1.7),
           Inches(6.6), Inches(5.0), labels, [("Total trays", qty, AMBER)],
           cat_font=11)
    # Metrics table.
    rows = []
    for m in trend.months:
        mm = m.metrics
        rows.append([m.label, _fmt(mm["total_qty"]), _fmt(mm["avg_per_day"]),
                     _fmt(mm["fuel_eff"]), f'{_fmt(mm["downtime_pct"])}%',
                     f'{_fmt(mm["repulp_rate"])}%'])
    _table(s, Inches(7.7), Inches(1.8), Inches(4.8),
           ["Month", "Trays", "Avg/d", "L/1k", "Down", "Rej"], rows,
           col_widths=[Inches(1.6), Inches(0.95), Inches(0.75), Inches(0.6),
                       Inches(0.5), Inches(0.4)],
           body_size=10, head_size=10)
    _caption(s, overall_summary)
    return s


def _improvement_plan_slide(prs, plan: list):
    """AI-written 'Production Improvement Plan' — numbered issues + actions."""
    s = _slide_header(prs, "Production Improvement Plan",
                      "AI-generated key issues & recommended actions")
    top = Inches(1.65)
    for i, item in enumerate(plan[:5]):
        title = str(item.get("title", "")).strip()
        detail = str(item.get("detail", "")).strip()
        action = str(item.get("action", "")).strip()
        if not title:
            continue
        _text(s, Inches(0.85), top, Inches(0.8), Inches(0.8),
              [(f"{i + 1:02d}", 22, True, AMBER)])
        runs = [(title, 15, True, WHITE)]
        if detail:
            runs.append((detail, 11, False, MUTE))
        if action:
            runs.append((f"Action: {action}", 11, False, AMBER))
        _text(s, Inches(1.7), top, Inches(10.8), Inches(1.0), runs)
        top += Inches(1.05)
    return s


def _downtime_comparison_slide(prs, trend: TrendData):
    s = _slide_header(prs, "Downtime Category Comparison",
                      " vs ".join(m.label.split(" ")[0] for m in trend.months))
    labels = [m.label.split(" ")[0][:3] for m in trend.months]
    # One series per cause, categories = months -> clustered columns.
    present = [c for c in CAUSE_ORDER if any(m.causes.get(c, 0) > 0 for m in trend.months)]
    series = []
    for i, c in enumerate(present):
        vals = [round(m.causes.get(c, 0) / 60, 1) for m in trend.months]
        series.append((c, vals, SERIES_COLORS[i % len(SERIES_COLORS)]))
    if series:
        _chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.85), Inches(1.7),
               Inches(7.4), Inches(5.0), labels, series, legend=True, cat_font=11)
    # Table: category x month (hours).
    headers = ["Category"] + [m.label.split(" ")[0][:3] for m in trend.months]
    rows = []
    for c in present:
        rows.append([c] + [_fmt(round(m.causes.get(c, 0) / 60, 1)) for m in trend.months])
    rows.append(["TOTAL"] + [_fmt(m.metrics["total_downtime_hrs"]) for m in trend.months])
    _table(s, Inches(8.45), Inches(1.9), Inches(4.05), headers, rows,
           body_size=10, head_size=10)
    return s


def _observations_slide(prs, trend: TrendData):
    """Auto-generated, data-driven observations across the period."""
    s = _slide_header(prs, "Observations & Trends",
                      "Auto-generated from the period's data")
    first, last = trend.months[0], trend.months[-1]
    obs = []

    def _delta(a, b):
        return None if not a else round((b - a) / a * 100)

    dq = _delta(first.metrics["total_qty"], last.metrics["total_qty"])
    if dq is not None:
        arrow = "↑" if dq > 0 else ("↓" if dq < 0 else "→")
        obs.append((f"Output {arrow} {abs(dq)}% over the period",
                    f"{first.label}: {_fmt(first.metrics['total_qty'])} trays  →  "
                    f"{last.label}: {_fmt(last.metrics['total_qty'])} trays."))
    de = _delta(first.metrics["fuel_eff"], last.metrics["fuel_eff"])
    if de is not None:
        better = "improved" if de < 0 else "worsened"
        obs.append((f"Fuel efficiency {better} ({de:+}%)",
                    f"{_fmt(first.metrics['fuel_eff'])} → {_fmt(last.metrics['fuel_eff'])} "
                    f"L per 1,000 trays (lower is better)."))
    dd = _delta(first.metrics["total_downtime_hrs"], last.metrics["total_downtime_hrs"])
    if dd is not None:
        arrow = "↑" if dd > 0 else "↓"
        obs.append((f"Downtime {arrow} {abs(dd)}%",
                    f"{_fmt(first.metrics['total_downtime_hrs'])} → "
                    f"{_fmt(last.metrics['total_downtime_hrs'])} hrs lost."))
    # Dominant downtime cause across the period.
    agg_causes = {}
    for m in trend.months:
        for c, v in m.causes.items():
            agg_causes[c] = agg_causes.get(c, 0) + v
    if agg_causes:
        top = max(agg_causes, key=agg_causes.get)
        share = round(agg_causes[top] / (sum(agg_causes.values()) or 1) * 100)
        obs.append((f"Top downtime cause: {top} ({share}% of lost time)",
                    "Focus improvement effort here for the biggest gain."))
    # Best / worst output month.
    best = max(trend.months, key=lambda m: m.metrics["total_qty"])
    worst = min(trend.months, key=lambda m: m.metrics["total_qty"])
    obs.append((f"Best month: {best.label} ({_fmt(best.metrics['total_qty'])} trays)",
                f"Lowest: {worst.label} ({_fmt(worst.metrics['total_qty'])} trays)."))

    top = Inches(1.7)
    for i, (head, detail) in enumerate(obs):
        num = f"{i + 1:02d}"
        _text(s, Inches(0.85), top, Inches(0.8), Inches(0.8),
              [(num, 22, True, AMBER)])
        _text(s, Inches(1.7), top, Inches(10.8), Inches(0.9),
              [(head, 15, True, WHITE), (detail, 11, False, MUTE)])
        top += Inches(0.92)
    return s


# --------------------------------------------------------------------------- #
# Entry point
# --------------------------------------------------------------------------- #
def build_report_pptx(
    db: Session,
    start: date | None = None,
    end: date | None = None,
    period_label: str = "Monthly",
) -> tuple[bytes, str]:
    trend = collect_trend_data(db, start, end)

    # One AI call for the whole deck (per-month one-liners, overall summary, and
    # the improvement plan), shared with the PDF/Excel reports. Returns {} when
    # AI is off or the call fails, so the deck always builds — AI text is additive.
    narrative = deck_narrative(trend)
    month_notes = narrative.get("months", {}) if narrative else {}

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    _title_slide(prs, trend)

    if not trend.months:
        _slide_header(prs, "No production data in this period",
                      "Pick a date range that contains shift entries.")
    else:
        # Per-month detail slides (with AI one-line summaries when available).
        for month in trend.months:
            notes = month_notes.get(month.key, {}) if isinstance(month_notes, dict) else {}
            _month_production_slide(prs, month, trend.targets, notes.get("production"))
            _month_downtime_slide(prs, month, notes.get("downtime"))
            _month_fuel_slide(prs, month, notes.get("fuel"))
        # Cross-month summary slides (only meaningful for 2+ months).
        if trend.multi:
            _trend_comparison_slide(prs, trend, narrative.get("overall_summary"))
            _downtime_comparison_slide(prs, trend)
            _observations_slide(prs, trend)
        # AI improvement-plan slide (only when AI produced one).
        plan = narrative.get("improvement_plan") if narrative else None
        if plan:
            _improvement_plan_slide(prs, plan)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), report_filename("pptx", start, end, period_label)
